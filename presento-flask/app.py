from flask import Flask, request, jsonify
from pptx import Presentation
from firebase_admin import credentials, initialize_app, storage
from flask_cors import CORS, cross_origin
import json
import subprocess
import numpy as np

# Initialize Firebase credentials and app
cred = credentials.Certificate("firebase_credentials.json")
initialize_app(cred, {"storageBucket": "presento-1d9cd.appspot.com"})

app = Flask(__name__)

cors = CORS(app)

@app.route("/api/makePPT", methods=["POST"])
@cross_origin()
def make():
    """
    Generate a PowerPoint presentation from the given JSON data and return the name of the created presentation along with the URLs of the uploaded files.

    Args:
        data (dict): A JSON object containing the necessary information to create the presentation. It should have the following keys:
            - "presentationTitle" (str): The title of the presentation.
            - "slide" (list): A list of slide objects. Each slide object should have the following keys:
                - "title" (str): The title of the slide.
                - "points" (list): A list of bullet points for the slide.

    Returns:
        dict: A JSON object containing the name of the created presentation and the URLs of the uploaded files. It has the following keys:
            - "name" (str): The name of the created presentation.
            - "url_pptx" (str): The URL of the uploaded PPTX file.
            - "url_pdf" (str): The URL of the uploaded PDF file.
    """
    def makePPT(data):
        def _add_leveled_bullet(_placeholder, _text, level=0):
            _prg = _placeholder.text_frame.add_paragraph()
            _prg.level = level
            _prg.text = _text

        # Create a presentation object
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = data["presentationTitle"]

        name = title.text.replace(" ", "_")
        name = name.replace("(", "")
        name = name.replace(")", "")
        

        for i in range(len(data["slide"])):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = data["slide"][i]["title"]
            for j in range(len(data["slide"][i]["points"])):
                _add_leveled_bullet(subtitle, data["slide"][i]["points"][j], 0)
        
        # Save the presentation to a file
        number = str(np.random.randint(low=1000))
        finalName = name + number
        prs.save(f"{finalName}.pptx")
        return finalName

    # Parse JSON data from request
    json_data = json.loads(request.get_json())
    print(json_data)
    finalName = makePPT(json_data)

    # Convert the created PPT to PDF using pptx2pdf command
    command = f"pptx2pdf {finalName}.pptx"
    subprocess.call(command, shell=True)

    file_path_pptx = f"{finalName}.pptx"
    file_path_pdf = f"{finalName}.pdf"

    # Convert the created PPT to PDF using pptx2pdf command
    bucket = storage.bucket()
    blob1 = bucket.blob(file_path_pptx)
    blob1.upload_from_filename(file_path_pptx)
    blob1.make_public()
    blob2 = bucket.blob(file_path_pdf)
    blob2.upload_from_filename(file_path_pdf)
    blob2.make_public()

    # Create a JSON response with the name and URLs of the uploaded files
    data = jsonify({"name": finalName, "url_pptx": blob1.public_url, "url_pdf": blob2.public_url})

    return data

# Run the Flask app
app.run(host="0.0.0.0", port=5000)